from copy import deepcopy
from pptx import Presentation
import os


class CertificateGenerator:

    def __init__(self, template_path):
        self.template_path = template_path

    def _replace_text(self, shape, replacements):
        if not hasattr(shape, "text_frame"):
            return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text

                for key, value in replacements.items():
                    text = text.replace(key, str(value))

                run.text = text

    def generate(
        self,
        output_path,
        name_ar,
        name_en,
        course_ar,
        course_en,
        date,
        hours
    ):

        prs = Presentation(self.template_path)

        replacements = {
            ">>الاسم<<": name_ar,
            ">>name<<": name_en,
            ">> [عنوان الدورة]<<": course_ar,
            ">>[عنوان الدورة]<<": course_ar,
            ">> [name of the course]<<": course_en,
            ">>[name of the course]<<": course_en,
            ">>التاريخ<<": date,
            ">>date<<": date,
            ">>عدد الساعات<<": hours,
            ">>hours<<": hours,
        }

        for slide in prs.slides:

            for shape in slide.shapes:

                self._replace_text(shape, replacements)

                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            self._replace_text(cell, replacements)

                if hasattr(shape, "shapes"):
                    for subshape in shape.shapes:
                        self._replace_text(subshape, replacements)

        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        prs.save(output_path)

        return output_path
